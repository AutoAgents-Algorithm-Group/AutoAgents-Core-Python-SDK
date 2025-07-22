from autoagentsai.client import ChatClient
import json
from autoagents_ai.utils import extract_json
import os
import requests
import tempfile
from typing import Optional, Union, List
from pptx import Presentation
from pptx.util import Inches

class create_ppt_agent:
    def __init__(self):
        pass

    def outline(self, prompt: str, file_path: Union[str, List[str]]):
        chat_client = ChatClient(
            agent_id="045c418f0dcf4adbb2f15031f06694d1",
            personal_auth_key="48cf18e0e0ca4b51bbf8fa60193ffb5c",
            personal_auth_secret="HWlQXZ5vxgrXDGEtTGGdsTFhJfr9rCmD",
            base_url="https://uat.agentspro.cn"
        )
        
        # 处理单文件或多文件情况
        if isinstance(file_path, str):
            files = [file_path]
        else:
            files = file_path
        
        print(f"Debug: 准备处理 {len(files)} 个文件: {files}")
        
        content = ""
        try:
            for event in chat_client.invoke(prompt, files=files):
                if event['type'] == 'start_bubble':
                    print(f"\n{'=' * 20} 消息气泡{event['bubble_id']}开始 {'=' * 20}")
                elif event['type'] == 'token':
                    print(event['content'], end='', flush=True)
                    content += event['content']
                elif event['type'] == 'end_bubble':
                    print(f"\n{'=' * 20} 消息气泡结束 {'=' * 20}")
                elif event['type'] == 'finish':
                    print(f"\n{'=' * 20} 对话完成 {'=' * 20}")
                    break
                elif event['type'] == 'error':
                    print(f"\nDebug: 收到错误事件: {event}")
                    break
                    
        except Exception as e:
            print(f"\nDebug: ChatClient.invoke 发生异常: {type(e).__name__}: {e}")
            # 如果流出现问题，返回错误信息而不是空字符串
            if not content.strip():
                content = f"Stream error: {str(e)}"
        
        print(f"\nDebug: 最终返回内容长度: {len(content)}")
        return content

    # def cover(self):
    #     pass

    # def content(self):
    #     pass

    # def conclusion(self):
    #     pass
    
    # def save(self, file_path: str):
    #     pass

    def download_image(self, url: str) -> Optional[str]:
        """下载远程图片到临时文件"""
        try:
            response = requests.get(url, stream=True)
            response.raise_for_status()
            
            # 创建临时文件
            suffix = '.jpg'  # 默认后缀
            if 'content-type' in response.headers:
                content_type = response.headers['content-type']
                if 'png' in content_type:
                    suffix = '.png'
                elif 'gif' in content_type:
                    suffix = '.gif'
                elif 'webp' in content_type:
                    suffix = '.webp'
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
            
            # 下载图片内容
            for chunk in response.iter_content(chunk_size=8192):
                temp_file.write(chunk)
            
            temp_file.close()
            return temp_file.name
            
        except Exception as e:
            print(f"下载图片失败: {url}, 错误: {e}")
            return None

    def fill(self, prompt: str, template_file_path: str = "template-1.pptx", output_file_path: str = "output-1.pptx"):
        # 创建一个PPT填充的agent
        chat_client = ChatClient(
            agent_id="868081a079604dd7ae57921c15f6297d",
            personal_auth_key="339859fa69934ea8b2b0ebd19d94d7f1",
            personal_auth_secret="93TsBecJplOawEipqAdF7TJ0g4IoBMtA",
            base_url="https://uat.agentspro.cn"
        )

        full_response = ""  # 用于存储完整的响应内容

        for event in chat_client.invoke(prompt):
            if event['type'] == 'start_bubble':
                print(f"\n{'=' * 20} 消息气泡{event['bubble_id']}开始 {'=' * 20}")
            elif event['type'] == 'reasoning_token':
                print(event['content'], end='', flush=True) 
            elif event['type'] == 'token':
                print(event['content'], end='', flush=True)
                full_response += event['content']
            elif event['type'] == 'end_bubble':
                print(f"\n{'=' * 20} 消息气泡结束 {'=' * 20}")
            elif event['type'] == 'finish':
                print(f"\n{'=' * 20} 对话完成 {'=' * 20}")
                break

        data = extract_json(full_response)

        print(data)

        # 加载 PPTX 模板
        prs = Presentation(template_file_path)

        # 用于存储需要清理的临时文件
        temp_files = []

        # 遍历所有幻灯片和形状
        for slide in prs.slides:
            # 处理文本替换
            for shape in slide.shapes:
                try:
                    # 检查是否有文本框
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        text_frame = getattr(shape, 'text_frame', None)
                        if text_frame:
                            for para in text_frame.paragraphs:
                                for run in para.runs:
                                    for key, val in data.items():
                                        # 只处理非图片的文本替换
                                        if not key.startswith('image'):
                                            placeholder = f"{{{{{key}}}}}"
                                            if placeholder in run.text:
                                                run.text = run.text.replace(placeholder, str(val))
                except Exception as e:
                    # 忽略无法处理的shape类型
                    continue
            
            # 处理图片替换 - 查找现有图片并替换
            pictures_to_replace = []
            image_keys = [key for key in data.keys() if key.startswith('image')]
            
            # 收集现有图片信息
            for shape in slide.shapes:
                try:
                    # 检查是否是图片形状
                    if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # 13 代表图片类型
                        pictures_to_replace.append({
                            'shape': shape,
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height
                        })
                except Exception as e:
                    continue
            
            # 替换图片
            for i, picture_info in enumerate(pictures_to_replace):
                # 按顺序匹配图片key（image1, image2, ...）
                if i < len(image_keys):
                    key = image_keys[i]
                    image_path = data[key]
                    
                    local_image_path = None
                    
                    # 判断是远程URL还是本地路径
                    if image_path.startswith(('http://', 'https://')):
                        # 下载远程图片
                        local_image_path = self.download_image(image_path)
                        if local_image_path:
                            temp_files.append(local_image_path)
                    elif os.path.exists(image_path):
                        # 本地文件路径
                        local_image_path = image_path
                    
                    # 替换图片
                    if local_image_path and os.path.exists(local_image_path):
                        try:
                            # 删除旧图片
                            old_shape = picture_info['shape']
                            sp = old_shape.element
                            sp.getparent().remove(sp)
                            
                            # 插入新图片到相同位置
                            slide.shapes.add_picture(
                                local_image_path, 
                                picture_info['left'], 
                                picture_info['top'], 
                                width=picture_info['width'], 
                                height=picture_info['height']
                            )
                            print(f"成功替换图片: {key}")
                        except Exception as e:
                            print(f"替换图片失败: {key}, 错误: {e}")
                    else:
                        print(f"警告: 无法获取图片: {image_path}")

        # 保存为新PPT
        prs.save(output_file_path)
        
        # 清理临时文件
        for temp_file in temp_files:
            try:
                os.unlink(temp_file)
                print(f"清理临时文件: {temp_file}")
            except Exception as e:
                print(f"清理临时文件失败: {temp_file}, 错误: {e}")