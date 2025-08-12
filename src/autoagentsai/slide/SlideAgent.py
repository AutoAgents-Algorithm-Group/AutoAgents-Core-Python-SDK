from ..client import ChatClient
from ..utils.extractor import extract_json
from ..utils.convertor import convert_csv_to_json_list
from .utils_slide import (
    download_image, 
    download_template,
    parse_markdown_text, 
    apply_inline_formatting, 
    enable_bullet, 
    fill_existing_table, 
    find_nearest_table,
    get_value_by_path,
    cleanup_temp_file,
    SimpleFileUploader,
    is_pure_placeholder,
    replace_mixed_placeholders
)
import os
import base64
import tempfile
from typing import List, Optional, Dict, Union
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN



class SlideAgent:
    def __init__(self):
        pass

    def outline(self, prompt: str, file_path_list: List[str]):
        chat_client = ChatClient(
            agent_id="045c418f0dcf4adbb2f15031f06694d1",
            personal_auth_key="48cf18e0e0ca4b51bbf8fa60193ffb5c",
            personal_auth_secret="HWlQXZ5vxgrXDGEtTGGdsTFhJfr9rCmD",
            base_url="https://uat.agentspro.cn"
        )
    
        print(f"Debug: 准备处理 {len(file_path_list)} 个文件: {file_path_list}")
        
        content = ""
        try:
            for event in chat_client.invoke(prompt, files=file_path_list):
                if event['type'] == 'start_bubble':
                    print(f"\n{'=' * 20} 消息气泡{event['bubble_id']}开始 {'=' * 20}")
                elif event['type'] == 'token':
                    print(event['content'], end='', flush=True)
                    content += event['content']
                elif event['type'] == 'end_bubble':
                    print(f"\n{'=' * 20} 消息气泡结束 {'=' * 20}")
                elif event['type'] == 'finish':
                    print(f"\n{'=' * 20} 对话完成 {'=' * 20}")
                    break
                elif event['type'] == 'error':
                    print(f"\nDebug: 收到错误事件: {event}")
                    break
                    
        except Exception as e:
            print(f"\nDebug: ChatClient.invoke 发生异常: {type(e).__name__}: {e}")
            # 如果流出现问题，返回错误信息而不是空字符串
            if not content.strip():
                content = f"Stream error: {str(e)}"
        
        print(f"\nDebug: 最终返回内容长度: {len(content)}")
        content = extract_json(content)
        return content


    def fill(self, 
             data: dict, 
             template_file_path: str, 
             output_file_path: Optional[str] = None,
             output_format: str = "local",
             personal_auth_key: Optional[str] = None,
             personal_auth_secret: Optional[str] = None,
             base_url: str = "https://uat.agentspro.cn") -> Union[str, Dict]:
        """
        使用数据填充PowerPoint模板
        
        Args:
            data: 要填充的数据字典
            template_file_path: 模板文件路径（支持本地路径和URL）
            output_file_path: 输出文件路径（当output_format为"local"时必需）
            output_format: 输出格式，支持 "local"、"base64"、"url"
            personal_auth_key: 当output_format为"url"时需要的认证密钥
            personal_auth_secret: 当output_format为"url"时需要的认证密钥
            base_url: 上传服务的基础URL
            
        Returns:
            str: 当output_format为"local"时返回文件路径，为"base64"时返回base64字符串
            Dict: 当output_format为"url"时返回上传结果字典
        """
        
        # 参数验证
        if output_format not in ["local", "base64", "url"]:
            raise ValueError(f"不支持的输出格式: {output_format}，支持的格式: local, base64, url")
        
        if output_format == "local" and not output_file_path:
            raise ValueError("当output_format为'local'时，必须提供output_file_path参数")
            
        if output_format == "url" and not personal_auth_key and not personal_auth_secret:
            raise ValueError("当output_format为'url'时，必须提供personal_auth_key和personal_auth_secret参数")
        # 用于存储需要清理的临时文件
        temp_files = []
        
        # 检查模板路径是否为URL，如果是则下载到临时文件
        actual_template_path = template_file_path
        is_template_from_url = False
        
        if template_file_path.startswith(('http://', 'https://')):
            print(f"检测到URL模板: {template_file_path}")
            downloaded_template = download_template(template_file_path)
            if downloaded_template:
                actual_template_path = downloaded_template
                temp_files.append(downloaded_template)
                is_template_from_url = True
                print(f"模板下载成功: {downloaded_template}")
            else:
                raise ValueError(f"无法下载模板文件: {template_file_path}")
        
        # 加载 PPTX 模板
        prs = Presentation(actual_template_path)

        # 处理远程图片下载
        processed_data = {}
        # 支持的图片文件后缀
        image_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.svg')
        
        def process_value(value):
            """递归处理数据值，支持CSV文件读取"""
            if isinstance(value, str):
                # 检查是否是CSV文件路径
                if value.endswith('.csv') and os.path.exists(value):
                    print(f"检测到CSV文件: {value}")
                    return convert_csv_to_json_list(value)
                # 检查是否是远程图片URL
                elif value.startswith(('http://', 'https://')):
                    # 检查URL是否以图片文件后缀结尾（忽略查询参数）
                    url_path = value.split('?')[0].lower()  # 去掉查询参数并转为小写
                    if url_path.endswith(image_extensions):
                        # 尝试下载远程图片
                        local_image_path = download_image(value)
                        if local_image_path:
                            temp_files.append(local_image_path)
                            print(f"成功下载图片: {value} -> {local_image_path}")
                            return local_image_path
                        else:
                            print(f"跳过下载失败的图片: {value}")
                            return None
                    else:
                        print(f"跳过非图片URL: {value} (不支持的文件类型)")
                        return value
                else:
                    return value
            elif isinstance(value, list):
                # 递归处理列表中的每个元素
                return [process_value(item) for item in value]
            elif isinstance(value, dict):
                # 递归处理字典中的每个值
                return {k: process_value(v) for k, v in value.items()}
            else:
                return value

        for key, value in data.items():
            processed_value = process_value(value)
            if processed_value is not None:
                processed_data[key] = processed_value

        # 1. 表格填充
        table_requests = []  # [(占位符形状, key, data)]
        all_tables = []  # 所有表格形状
        
        print(f"开始扫描PPT模板中的占位符...")
        for slide_idx, slide in enumerate(prs.slides):
            print(f"扫描第 {slide_idx + 1} 页...")
            for shape in slide.shapes:
                # 收集表格占位符
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text.startswith("{{") and text.endswith("}}"):
                        print(f"  找到占位符: {text}")
                    if text.startswith("{{#") and text.endswith("}}"):
                        path = text[3:-2].strip()  # 去掉 {{# 和 }}
                        print(f"找到表格占位符: {{#{path}}}")
                        table_data = get_value_by_path(processed_data, path)
                        if table_data is not None:
                            print(f"表格占位符 {{#{path}}} 数据解析成功，{len(table_data) if isinstance(table_data, list) else 1} 条记录")
                            table_requests.append((shape, path, table_data))
                        else:
                            print(f"表格占位符 {{#{path}}} 数据解析失败")
                
                # 收集所有表格
                if shape.has_table:
                    all_tables.append(shape)
        
        # 为每个表格占位符找到最近的表格并填充
        shapes_to_remove = []
        processed_tables = set()
        
        for placeholder_shape, path, table_data in table_requests:
            available_tables = [t for t in all_tables if id(t) not in processed_tables]
            if not available_tables:
                available_tables = all_tables
            
            nearest_table_shape = find_nearest_table(placeholder_shape, available_tables)
            if nearest_table_shape:
                print(f"占位符 '{{#{path}}}' 匹配到最近的表格")
                fill_existing_table(nearest_table_shape.table, table_data)
                processed_tables.add(id(nearest_table_shape))
            
            shapes_to_remove.append(placeholder_shape)
        
        # 删除表格占位符文本框
        for shape in shapes_to_remove:
            shape._element.getparent().remove(shape._element)
        
        # 2. 文本、图片填充
        for slide in prs.slides:
            for shape in list(slide.shapes):  # list() to allow removal
                if not shape.has_text_frame:
                    continue
            
                text = shape.text.strip()
                
                # 检查是否包含占位符
                if "{{" in text and "}}" in text:
                    # 检查是否为纯占位符
                    pure_placeholder = is_pure_placeholder(text)
                    
                    if pure_placeholder:
                        # 纯占位符模式（原有逻辑）
                        path = pure_placeholder
                        content_type = "text"

                        # 判断类型前缀
                        if path.startswith("@"):
                            path = path[1:]
                            content_type = "image"
                        elif path.startswith("#"):
                            # 表格已经在上面处理过了，跳过
                            continue

                        value = get_value_by_path(processed_data, path)
                        if value is None:
                            continue

                        if content_type == "text":
                            # 检查是否包含Markdown格式
                            if isinstance(value, str) and any(marker in value for marker in ['*', '#', '`', '\n']):
                                # 使用Markdown解析
                                parse_markdown_text(shape.text_frame, value)
                            elif isinstance(value, list):
                                # 处理列表数据，每项作为bullet point
                                tf = shape.text_frame
                                tf.clear()
                                for i, item in enumerate(value):
                                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                                    if isinstance(item, str) and any(marker in item for marker in ['*', '#', '`']):
                                        apply_inline_formatting(p, item)
                                    else:
                                        p.text = str(item)
                                        p.font.size = Pt(14)
                                        p.alignment = PP_ALIGN.LEFT
                                        enable_bullet(p)
                            else:
                                # 普通文本
                                shape.text_frame.text = str(value)

                        elif content_type == "image":
                            # 获取位置并删除原文本框
                            left, top, width, height = shape.left, shape.top, shape.width, shape.height
                            slide.shapes._spTree.remove(shape._element)
                                
                            # 确保图片路径存在
                            if os.path.exists(value):
                                slide.shapes.add_picture(value, left, top, width=width, height=height)
                                print(f"成功替换图片: {path}")
                            else:
                                print(f"警告: 图片文件不存在: {value}")
                    
                    else:
                        # 混合文本模式（新功能）
                        replaced_text = replace_mixed_placeholders(text, processed_data)
                        shape.text_frame.text = replaced_text
                        print(f"混合文本替换: '{text}' -> '{replaced_text}'")

        # 根据输出格式处理结果
        result = None
        temp_output_path = None
        
        try:
            if output_format == "local":
                # 直接保存到指定路径
                prs.save(output_file_path)
                print(f"✅ PPT已保存到: {output_file_path}")
                result = output_file_path
                
            elif output_format == "base64":
                # 保存到临时文件，然后转换为base64
                temp_fd, temp_output_path = tempfile.mkstemp(suffix='.pptx')
                os.close(temp_fd)
                temp_files.append(temp_output_path)
                
                prs.save(temp_output_path)
                
                # 读取文件并转换为base64
                with open(temp_output_path, 'rb') as f:
                    file_bytes = f.read()
                    base64_str = base64.b64encode(file_bytes).decode('utf-8')
                
                print(f"✅ PPT已转换为base64格式 (大小: {len(base64_str)} 字符)")
                result = base64_str
                
            elif output_format == "url":
                # 保存到临时文件，然后上传
                temp_fd, temp_output_path = tempfile.mkstemp(suffix='.pptx')
                os.close(temp_fd)
                temp_files.append(temp_output_path)
                
                prs.save(temp_output_path)
                
                # 创建上传器并上传文件
                uploader = SimpleFileUploader(personal_auth_key, personal_auth_secret, base_url)
                
                # 生成文件名
                filename = f"filled_presentation_{os.path.basename(temp_output_path)}"
                
                with open(temp_output_path, 'rb') as f:
                    upload_result = uploader.upload(f, filename)
                
                if upload_result.get("success"):
                    print(f"✅ PPT已上传成功，文件ID: {base_url}/api/fs/{upload_result['fileId']}")
                    result = {
                        "fileId": upload_result['fileId'],
                        "fileUrl": f"{base_url}/api/fs/{upload_result['fileId']}"
                    }
                else:
                    raise Exception(f"文件上传失败: {upload_result.get('error', '未知错误')}")
            
        finally:
            # 清理临时文件（包括下载的图片、模板文件和输出临时文件）
            for temp_file in temp_files:
                try:
                    cleanup_temp_file(temp_file)
                    print(f"清理临时文件: {temp_file}")
                except Exception as e:
                    print(f"清理临时文件失败: {temp_file}, 错误: {e}")
        
        return result