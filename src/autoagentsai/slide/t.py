class TextStylePreserver:
    """完整的文本样式保存和恢复工具类"""
    
    @staticmethod
    def capture_complete_style(paragraph):
        """捕获段落的完整样式信息"""
        style_info = {
            'paragraph_level': {},
            'run_level': [],
            'text_frame_level': {}
        }
        
        # 段落级别样式
        style_info['paragraph_level'] = {
            'alignment': paragraph.alignment,
            'level': paragraph.level,
            'space_before': paragraph.space_before,
            'space_after': paragraph.space_after,
            'line_spacing': paragraph.line_spacing,
            'has_bullet': False,
            'bullet_char': None
        }
        
        # 检查是否有项目符号
        try:
            p_element = paragraph._p
            pPr = p_element.pPr
            if pPr is not None:
                # 查找项目符号字符
                buChar_elements = pPr.xpath('.//a:buChar', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if buChar_elements:
                    style_info['paragraph_level']['has_bullet'] = True
                    char_attr = buChar_elements[0].get('char')
                    if char_attr:
                        style_info['paragraph_level']['bullet_char'] = char_attr
                    else:
                        style_info['paragraph_level']['bullet_char'] = "•"  # 默认符号
                # 检查其他类型的项目符号（如数字编号等）
                elif pPr.xpath('.//a:buFont | .//a:buAutoNum | .//a:buBlip', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                    style_info['paragraph_level']['has_bullet'] = True
                    style_info['paragraph_level']['bullet_char'] = "•"  # 默认符号
        except Exception:
            pass
        
        # 获取段落的XML以保存更多属性
        try:
            # 保存段落属性的XML片段
            style_info['paragraph_level']['xml_props'] = p_element.xml if hasattr(p_element, 'xml') else None
        except:
            pass
        
        # Run级别样式（每个run）
        for run in paragraph.runs:
            # 处理可能为None的属性
            font_bold = run.font.bold
            if font_bold is None:
                font_bold = False  # 默认不bold
            
            font_italic = run.font.italic  
            if font_italic is None:
                font_italic = False  # 默认不italic
                
            font_underline = run.font.underline
            if font_underline is None:
                font_underline = False  # 默认无下划线
            
            run_style = {
                'text': run.text,
                'font_name': run.font.name,
                'font_size': run.font.size,
                'font_bold': font_bold,
                'font_italic': font_italic,
                'font_underline': font_underline,
                'font_color_rgb': None,
                'font_color_theme': None,
                'hyperlink': None,
                'font_element_xml': None  # 保存字体元素的XML
            }
            
            # 字体颜色处理
            try:
                if run.font.color.rgb:
                    run_style['font_color_rgb'] = run.font.color.rgb
                elif run.font.color.theme_color:
                    run_style['font_color_theme'] = run.font.color.theme_color
            except AttributeError:
                pass
            
            # 尝试保存字体的原始XML以确保完整性
            try:
                if hasattr(run, '_r'):
                    run_style['font_element_xml'] = run._r.xml
            except:
                pass
            
            # 超链接处理
            try:
                if hasattr(run, '_r') and run._r.get('hlinkClick'):
                    run_style['hyperlink'] = run._r.get('hlinkClick')
            except:
                pass
                
            style_info['run_level'].append(run_style)
        
        return style_info
    
    @staticmethod
    def capture_text_frame_style(text_frame):
        """捕获文本框级别的样式"""
        return {
            'vertical_anchor': text_frame.vertical_anchor,
            'margin_left': text_frame.margin_left,
            'margin_right': text_frame.margin_right,
            'margin_top': text_frame.margin_top,
            'margin_bottom': text_frame.margin_bottom,
            'word_wrap': text_frame.word_wrap,
            'auto_size': text_frame.auto_size
        }
    
    @staticmethod
    def apply_style_to_new_text(paragraph, style_info, new_text):
        """将保存的样式应用到新文本上，完全保留格式"""
        
        # 应用段落级别样式
        para_style = style_info['paragraph_level']
        if para_style.get('alignment') is not None:
            paragraph.alignment = para_style['alignment']
        if para_style.get('level') is not None:
            paragraph.level = para_style['level']
        if para_style.get('space_before') is not None:
            paragraph.space_before = para_style['space_before']
        if para_style.get('space_after') is not None:
            paragraph.space_after = para_style['space_after']
        if para_style.get('line_spacing') is not None:
            paragraph.line_spacing = para_style['line_spacing']
        
        # 恢复项目符号
        if para_style.get('has_bullet', False):
            bullet_char = para_style.get('bullet_char', "•")
            enable_bullet(paragraph, bullet_char)
        
        # 保留第一个run的样式信息，然后替换文本
        if paragraph.runs and style_info['run_level']:
            original_run_style = style_info['run_level'][0]
            
            # 清除所有现有runs（除了第一个）
            while len(paragraph.runs) > 1:
                run = paragraph.runs[-1]
                run._r.getparent().remove(run._r)
            
            # 获取第一个run并设置新文本
            if paragraph.runs:
                run = paragraph.runs[0]
                run.text = new_text
                
                # 强化字体样式应用 - 确保字体名称被正确设置
                try:
                    if original_run_style.get('font_name'):
                        original_font_name = original_run_style['font_name']
                        run.font.name = original_font_name
                        
                        # 验证字体是否设置成功，如果不匹配则尝试XML级别设置
                        if run.font.name != original_font_name:
                            try:
                                # 直接操作底层XML元素
                                if hasattr(run, '_r'):
                                    r_element = run._r
                                    # 获取或创建rPr元素
                                    rPr = r_element.get_or_add_rPr()
                                    
                                    # 移除现有的latin字体设置（如果有）
                                    for latin in rPr.xpath('.//a:latin', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                                        rPr.remove(latin)
                                    
                                    # 添加新的latin字体设置
                                    from pptx.oxml import parse_xml
                                    latin_font = parse_xml(f'<a:latin xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{original_font_name}"/>')
                                    rPr.insert(0, latin_font)
                            except Exception:
                                pass  # 字体设置失败时静默处理
                except Exception:
                    pass  # 字体名称设置异常时静默处理
                
                try:
                    if original_run_style.get('font_size'):
                        run.font.size = original_run_style['font_size']
                except Exception as e:
                    print(f"⚠️ 字体大小设置异常: {e}")
                
                # 明确设置布尔属性
                try:
                    run.font.bold = original_run_style.get('font_bold', False)
                except:
                    pass
                try:
                    run.font.italic = original_run_style.get('font_italic', False)
                except:
                    pass
                try:
                    run.font.underline = original_run_style.get('font_underline', False)
                except:
                    pass
                
                # 应用字体颜色
                try:
                    if original_run_style.get('font_color_rgb'):
                        run.font.color.rgb = original_run_style['font_color_rgb']
                    elif original_run_style.get('font_color_theme'):
                        run.font.color.theme_color = original_run_style['font_color_theme']
                except Exception as e:
                    print(f"⚠️ 字体颜色设置异常: {e}")
            else:
                # 如果没有现有runs，创建一个新的run并应用样式
                paragraph.text = new_text
                if paragraph.runs:
                    run = paragraph.runs[0]
                    # 强化字体样式应用（备用情况）
                    try:
                        if original_run_style.get('font_name'):
                            original_font_name = original_run_style['font_name']
                            run.font.name = original_font_name
                            # 如果字体名称不匹配，尝试XML级别设置
                            if run.font.name != original_font_name:
                                try:
                                    # 直接操作底层XML元素
                                    if hasattr(run, '_r'):
                                        r_element = run._r
                                        # 获取或创建rPr元素
                                        rPr = r_element.get_or_add_rPr()
                                        
                                        # 移除现有的latin字体设置（如果有）
                                        for latin in rPr.xpath('.//a:latin', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                                            rPr.remove(latin)
                                        
                                        # 添加新的latin字体设置
                                        from pptx.oxml import parse_xml
                                        latin_font = parse_xml(f'<a:latin xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{original_font_name}"/>')
                                        rPr.insert(0, latin_font)
                                except Exception:
                                    pass  # 字体设置失败时静默处理
                    except Exception:
                        pass  # 字体名称设置异常时静默处理
                    
                    try:
                        if original_run_style.get('font_size'):
                            run.font.size = original_run_style['font_size']
                    except:
                        pass
                    try:
                        run.font.bold = original_run_style.get('font_bold', False)
                    except:
                        pass
                    try:
                        run.font.italic = original_run_style.get('font_italic', False)
                    except:
                        pass
                    try:
                        run.font.underline = original_run_style.get('font_underline', False)
                    except:
                        pass
                    try:
                        if original_run_style.get('font_color_rgb'):
                            run.font.color.rgb = original_run_style['font_color_rgb']
                        elif original_run_style.get('font_color_theme'):
                            run.font.color.theme_color = original_run_style['font_color_theme']
                    except:
                        pass
        else:
            # 如果没有样式信息，使用简单的文本替换
            paragraph.text = new_text
    
    @staticmethod
    def apply_text_frame_style(text_frame, style_info):
        """应用文本框级别样式"""
        if style_info.get('vertical_anchor') is not None:
            text_frame.vertical_anchor = style_info['vertical_anchor']
        if style_info.get('margin_left') is not None:
            text_frame.margin_left = style_info['margin_left']
        if style_info.get('margin_right') is not None:
            text_frame.margin_right = style_info['margin_right']
        if style_info.get('margin_top') is not None:
            text_frame.margin_top = style_info['margin_top']
        if style_info.get('margin_bottom') is not None:
            text_frame.margin_bottom = style_info['margin_bottom']
        if style_info.get('word_wrap') is not None:
            text_frame.word_wrap = style_info['word_wrap']
        if style_info.get('auto_size') is not None:
            text_frame.auto_size = style_info['auto_size']


def replace_text_preserve_format(text_frame, new_text):
    """替换文本并完全保留格式的核心函数"""
    if not text_frame.paragraphs:
        return
    
    # 捕获文本框级别样式
    tf_style = TextStylePreserver.capture_text_frame_style(text_frame)
    
    # 检查新文本是否包含多行
    new_lines = new_text.split('\n')
    original_paragraphs = list(text_frame.paragraphs)
    
    # 捕获所有现有段落的样式
    para_styles = []
    for para in original_paragraphs:
        para_styles.append(TextStylePreserver.capture_complete_style(para))
    
    # 如果新文本只有一行，使用原来的简单逻辑
    if len(new_lines) == 1:
        # 捕获第一个段落的样式作为模板
        first_para = text_frame.paragraphs[0]
        para_style = para_styles[0] if para_styles else TextStylePreserver.capture_complete_style(first_para)
        
        # 删除所有现有段落（除了第一个）
        while len(text_frame.paragraphs) > 1:
            p = text_frame.paragraphs[-1]
            text_frame._element.remove(p._p)
        
        # 在第一个段落应用新文本和样式
        TextStylePreserver.apply_style_to_new_text(first_para, para_style, new_text)
    else:
        # 多行文本，需要保留每行的段落格式
        # 首先处理第一行
        if para_styles:
            first_para = text_frame.paragraphs[0]
            TextStylePreserver.apply_style_to_new_text(first_para, para_styles[0], new_lines[0])
        
        # 处理剩余行
        for i, line in enumerate(new_lines[1:], 1):
            # 决定使用哪个样式模板
            if i < len(para_styles):
                # 使用对应的原始段落样式
                style_template = para_styles[i]
            elif para_styles:
                # 使用最后一个段落的样式作为模板
                style_template = para_styles[-1]
            else:
                # 使用第一个段落的样式作为模板
                style_template = para_styles[0] if para_styles else TextStylePreserver.capture_complete_style(text_frame.paragraphs[0])
            
            # 添加新段落或使用现有段落
            if i < len(text_frame.paragraphs):
                # 使用现有段落
                para = text_frame.paragraphs[i]
            else:
                # 添加新段落
                para = text_frame.add_paragraph()
            
            # 应用样式和文本
            TextStylePreserver.apply_style_to_new_text(para, style_template, line)
        
        # 删除多余的段落（如果新文本行数少于原段落数）
        while len(text_frame.paragraphs) > len(new_lines):
            p = text_frame.paragraphs[-1]
            text_frame._element.remove(p._p)
    
    # 恢复文本框级别样式
    TextStylePreserver.apply_text_frame_style(text_frame, tf_style)


def process_list_preserve_format(text_frame, list_data):
    """处理列表数据并完全保留格式"""
    if not text_frame.paragraphs or not list_data:
        return
        
    # 捕获文本框级别样式
    tf_style = TextStylePreserver.capture_text_frame_style(text_frame)
    
    # 捕获第一个段落的样式作为模板
    first_para = text_frame.paragraphs[0]
    para_style = TextStylePreserver.capture_complete_style(first_para)
    
    # 删除所有现有段落（除了第一个）
    while len(text_frame.paragraphs) > 1:
        p = text_frame.paragraphs[-1]
        text_frame._element.remove(p._p)
    
    # 处理第一个列表项
    if list_data:
        first_item = str(list_data[0])
        TextStylePreserver.apply_style_to_new_text(first_para, para_style, first_item)
        enable_bullet(first_para)
        
        # 处理剩余的列表项
        for item in list_data[1:]:
            new_para = text_frame.add_paragraph()
            TextStylePreserver.apply_style_to_new_text(new_para, para_style, str(item))
            enable_bullet(new_para)
    
    # 恢复文本框级别样式
    TextStylePreserver.apply_text_frame_style(text_frame, tf_style)