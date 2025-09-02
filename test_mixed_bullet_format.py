#!/usr/bin/env python3
"""
专门测试混合文本中项目符号保留的脚本
"""
import sys
import os
from pptx import Presentation
from pptx.util import Pt

# 添加项目路径
sys.path.append(os.path.join(os.path.dirname(__file__), '.'))

from src.autoagentsai.slide.pptx2pptx_agent import enable_bullet, PPTX2PPTXAgent

def create_mixed_bullet_template():
    """创建包含混合文本项目符号的测试模板"""
    prs = Presentation()
    
    # 创建幻灯片
    slide_layout = prs.slide_layouts[1]  # 使用标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title = slide.shapes.title
    title.text = "混合文本项目符号测试"
    
    # 添加包含混合文本项目符号的文本框
    left = Pt(50)
    top = Pt(100)
    width = Pt(500)
    height = Pt(300)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    # 第一个段落 - 混合文本（类似用户的情况）
    p1 = text_frame.paragraphs[0]
    p1.text = "输入（Input）：组件托盘号：{{cover.title}}"
    enable_bullet(p1, "◆")  # 使用菱形符号
    
    # 第二个段落 - 混合文本  
    p2 = text_frame.add_paragraph()
    p2.text = "输出（Output）：组件序列号：{{serial_number}}"
    enable_bullet(p2, "◆")
    
    # 第三个段落 - 固定文本
    p3 = text_frame.add_paragraph()
    p3.text = "状态：已完成"
    enable_bullet(p3, "◆")
    
    # 保存模板
    template_path = "playground/slide/input/mixed_bullet_template.pptx"
    prs.save(template_path)
    print(f"✅ 创建混合项目符号测试模板: {template_path}")
    return template_path

def test_mixed_bullet_preservation():
    """测试混合文本项目符号保留功能"""
    # 创建测试模板
    template_path = create_mixed_bullet_template()
    
    agent = PPTX2PPTXAgent()
    
    # 测试数据
    test_data = {
        "cover": {
            "title": "新的组件托盘号123"
        },
        "serial_number": "ABC-789-XYZ"
    }
    
    try:
        print("🧪 开始测试混合文本项目符号保留...")
        
        result = agent.fill(
            data=test_data,
            template_file_path=template_path,
            output_file_path="playground/slide/output/mixed_bullet_test_output.pptx",
            verbose=True
        )
        
        print(f"✅ 测试完成，输出文件: {result}")
        print("\n🔍 请检查输出的PPT文件，确认：")
        print("1. 项目符号 ◆ 是否保留")
        print("2. 文本替换是否正确")
        print("3. 格式是否完整保持")
        return True
        
    except Exception as e:
        print(f"❌ 测试失败: {e}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    # 确保目录存在
    os.makedirs("playground/slide/input", exist_ok=True)
    os.makedirs("playground/slide/output", exist_ok=True)
    
    success = test_mixed_bullet_preservation()
    if success:
        print("\n🎉 混合文本项目符号测试完成！")
    else:
        print("\n💥 混合文本项目符号测试失败！")
        sys.exit(1)
