#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
多数据填充功能测试文件
演示如何使用同一模板页面填入不同内容
"""

import os
import sys
sys.path.append(os.path.join(os.path.dirname(__file__), '../../'))

from src.autoagentsai.slide.pptx2pptx_agent import PPTX2PPTXAgent

def create_students_data():
    """创建学生成绩单数据 - 多个学生，同一个模板页面"""
    return [
        {
            "name": "张三",
            "student_id": "2024001", 
            "score": 95,
            "grade": "A",
            "class": "计算机科学1班",
            "photo": "https://via.placeholder.com/150x200/4CAF50/FFFFFF?text=张三",
            "subjects": [
                {"name": "数学", "score": 98, "teacher": "王老师"},
                {"name": "语文", "score": 92, "teacher": "李老师"},
                {"name": "英语", "score": 95, "teacher": "陈老师"}
            ],
            "achievements": [
                "数学竞赛第一名",
                "优秀学生干部",
                "编程大赛冠军"
            ],
            "comments": "该学生表现优异，思维敏捷，是班级的榜样。"
        },
        {
            "name": "李四",
            "student_id": "2024002",
            "score": 88,
            "grade": "B",
            "class": "计算机科学1班", 
            "photo": "https://via.placeholder.com/150x200/2196F3/FFFFFF?text=李四",
            "subjects": [
                {"name": "数学", "score": 85, "teacher": "王老师"},
                {"name": "语文", "score": 91, "teacher": "李老师"},
                {"name": "英语", "score": 88, "teacher": "陈老师"}
            ],
            "achievements": [
                "语文作文比赛第二名",
                "团队合作优秀奖",
                "社会实践积极分子"
            ],
            "comments": "该学生学习认真，团队协作能力强，有很大潜力。"
        },
        {
            "name": "王五",
            "student_id": "2024003",
            "score": 92,
            "grade": "A",
            "class": "计算机科学1班",
            "photo": "https://via.placeholder.com/150x200/FF9800/FFFFFF?text=王五", 
            "subjects": [
                {"name": "数学", "score": 94, "teacher": "王老师"},
                {"name": "语文", "score": 89, "teacher": "李老师"},
                {"name": "英语", "score": 93, "teacher": "陈老师"}
            ],
            "achievements": [
                "科技创新大赛一等奖",
                "最佳创意奖",
                "学习进步奖"
            ],
            "comments": "该学生富有创新精神，动手能力强，是未来的科技人才。"
        }
    ]

def create_products_data():
    """创建产品介绍数据 - 多个产品，同一个模板页面"""
    return [
        {
            "product_name": "智能助手 Pro 3.0",
            "product_code": "AI-PRO-300",
            "price": "¥299/月",
            "category": "人工智能",
            "image": "https://via.placeholder.com/400x300/4CAF50/FFFFFF?text=AI+PRO",
            "features": [
                "自然语言理解",
                "多模态交互", 
                "智能任务规划",
                "企业级安全"
            ],
            "specifications": [
                {"item": "CPU要求", "value": "4核以上"},
                {"item": "内存要求", "value": "8GB以上"},
                {"item": "存储空间", "value": "50GB"}
            ],
            "description": "领先的人工智能助手，为企业提供全方位的智能化解决方案。"
        },
        {
            "product_name": "数据分析大师",
            "product_code": "DA-MASTER-200",
            "price": "¥199/月",
            "category": "数据分析",
            "image": "https://via.placeholder.com/400x300/2196F3/FFFFFF?text=DATA+MASTER",
            "features": [
                "可视化图表",
                "实时数据处理",
                "智能报告生成",
                "多数据源集成"
            ],
            "specifications": [
                {"item": "数据源", "value": "支持100+种"},
                {"item": "并发用户", "value": "1000+"},
                {"item": "报告模板", "value": "500+种"}
            ],
            "description": "专业的数据分析工具，帮助企业从海量数据中挖掘价值。"
        },
        {
            "product_name": "云端存储专家",
            "product_code": "CLOUD-STORE-100", 
            "price": "¥99/月",
            "category": "云存储",
            "image": "https://via.placeholder.com/400x300/FF9800/FFFFFF?text=CLOUD+STORE",
            "features": [
                "无限容量",
                "多端同步",
                "企业级加密",
                "版本管理"
            ],
            "specifications": [
                {"item": "上传速度", "value": "100MB/s"},
                {"item": "下载速度", "value": "200MB/s"},
                {"item": "可用性", "value": "99.99%"}
            ],
            "description": "安全可靠的云端存储服务，为企业数据保驾护航。"
        }
    ]

def test_students_report():
    """测试学生成绩单场景"""
    print("🎓 测试场景1: 学生成绩单生成")
    print("-" * 50)
    
    agent = PPTX2PPTXAgent()
    students_data = create_students_data()
    
    try:
        result = agent.fill_multiple_data(
            data_list=students_data,
            template_file_path="playground/slide/input/test_template.pptx",
            template_slide_index=0,  # 使用第1页作为学生成绩单模板
            output_file_path="playground/slide/output/students_report.pptx",
            verbose=True
        )
        print(f"✅ 学生成绩单生成成功: {result}")
        return True
    except Exception as e:
        print(f"❌ 学生成绩单生成失败: {e}")
        return False

def test_products_catalog():
    """测试产品目录场景"""
    print("\n📦 测试场景2: 产品目录生成")
    print("-" * 50)
    
    agent = PPTX2PPTXAgent()
    products_data = create_products_data()
    
    try:
        result = agent.fill_multiple_data(
            data_list=products_data,
            template_file_path="playground/slide/input/test_template.pptx",
            template_slide_index=1,  # 使用第2页作为产品介绍模板（如果有的话）
            output_file_path="playground/slide/output/products_catalog.pptx",
            verbose=True
        )
        print(f"✅ 产品目录生成成功: {result}")
        return True
    except Exception as e:
        print(f"❌ 产品目录生成失败: {e}")
        # 如果第2页不存在，尝试使用第1页
        print("📋 尝试使用第1页作为模板...")
        try:
            result = agent.fill_multiple_data(
                data_list=products_data,
                template_file_path="playground/slide/input/test_template.pptx",
                template_slide_index=0,  # 使用第1页作为模板
                output_file_path="playground/slide/output/products_catalog_alt.pptx",
                verbose=True
            )
            print(f"✅ 产品目录生成成功（使用第1页模板): {result}")
            return True
        except Exception as e2:
            print(f"❌ 产品目录生成仍然失败: {e2}")
            return False

def test_base64_output():
    """测试Base64输出格式"""
    print("\n💾 测试场景3: Base64输出格式")
    print("-" * 50)
    
    agent = PPTX2PPTXAgent()
    # 使用较少的数据项以减少输出大小
    simple_data = create_students_data()[:2]  # 只使用前两个学生
    
    try:
        result = agent.fill_multiple_data(
            data_list=simple_data,
            template_file_path="playground/slide/input/test_template.pptx",
            template_slide_index=0,
            output_format="base64",
            verbose=False  # 减少输出噪音
        )
        
        # 保存base64到文件
        with open("playground/slide/output/multiple_data_base64.txt", "w") as f:
            f.write(result)
            
        print(f"✅ Base64输出成功，长度: {len(result)} 字符")
        print(f"📁 Base64内容已保存到: playground/slide/output/multiple_data_base64.txt")
        return True
    except Exception as e:
        print(f"❌ Base64输出失败: {e}")
        return False

def create_custom_template_demo():
    """创建一个简单的自定义模板演示（如果原模板不存在）"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    template_path = "playground/slide/input/demo_template.pptx"
    
    # 如果演示模板不存在，创建一个简单的
    if not os.path.exists(template_path):
        print("📝 创建演示模板...")
        
        prs = Presentation()
        
        # 删除默认幻灯片
        if len(prs.slides) > 0:
            slide_id = prs.slides[0].slide_id
            prs.part.drop_rel(prs.slides._sldIdLst[0].rId)
            del prs.slides._sldIdLst[0]
        
        # 添加一张演示幻灯片
        layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "{{name}} - 个人信息"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        
        # 添加照片占位符
        photo_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(2))
        photo_frame = photo_box.text_frame
        photo_frame.text = "{{@photo}}"
        
        # 添加基本信息
        info_box = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(4), Inches(3))
        info_frame = info_box.text_frame
        info_frame.text = """学号: {{student_id}}
班级: {{class}}
总分: {{score}}
等级: {{grade}}

评语: {{comments}}"""
        
        # 添加成就列表
        achievements_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(7), Inches(2))
        achievements_frame = achievements_box.text_frame
        achievements_frame.text = "主要成就:\n{{achievements}}"
        
        # 保存模板
        os.makedirs(os.path.dirname(template_path), exist_ok=True)
        prs.save(template_path)
        print(f"✅ 演示模板创建成功: {template_path}")
        
        return template_path
    
    return template_path

def test_with_custom_template():
    """使用自定义模板测试"""
    print("\n🎨 测试场景4: 自定义模板演示")
    print("-" * 50)
    
    # 创建或获取演示模板
    template_path = create_custom_template_demo()
    
    agent = PPTX2PPTXAgent()
    students_data = create_students_data()
    
    try:
        result = agent.fill_multiple_data(
            data_list=students_data,
            template_file_path=template_path,
            template_slide_index=0,
            output_file_path="playground/slide/output/custom_template_result.pptx",
            verbose=True
        )
        print(f"✅ 自定义模板填充成功: {result}")
        return True
    except Exception as e:
        print(f"❌ 自定义模板填充失败: {e}")
        return False

def main():
    """主测试函数"""
    print("🚀 多数据填充功能测试")
    print("=" * 60)
    print("💡 功能说明：使用同一个模板页面，填入不同内容生成多页PPT")
    print("=" * 60)
    
    # 确保输出目录存在
    os.makedirs("playground/slide/output", exist_ok=True)
    
    # 运行测试
    results = []
    test_functions = [
        ("学生成绩单", test_students_report),
        ("产品目录", test_products_catalog),
        ("Base64输出", test_base64_output),
        ("自定义模板", test_with_custom_template)
    ]
    
    for test_name, test_func in test_functions:
        try:
            success = test_func()
            results.append((test_name, success))
        except Exception as e:
            print(f"❌ {test_name}测试出现异常: {e}")
            results.append((test_name, False))
    
    # 测试结果总结
    print("\n📊 测试结果总结")
    print("=" * 60)
    passed = sum(1 for _, success in results if success)
    total = len(results)
    
    for test_name, success in results:
        status = "✅ 通过" if success else "❌ 失败"
        print(f"{test_name}: {status}")
    
    print(f"\n🎯 总计: {passed}/{total} 个测试通过")
    
    if passed == total:
        print("🎉 所有测试通过！")
        print("\n💡 使用场景总结：")
        print("• 📚 批量生成学生成绩单/报告卡")
        print("• 📦 批量生成产品介绍页面")
        print("• 📄 批量生成员工简历/档案")
        print("• 📊 批量生成数据报告页面")
        print("• 🎫 批量生成证书/奖状")
    else:
        print("⚠️ 部分测试失败，请检查错误信息")
    
    print("\n📁 输出文件:")
    if os.path.exists("playground/slide/output"):
        for file in os.listdir("playground/slide/output"):
            if file.endswith(('.pptx', '.txt')):
                print(f"  • playground/slide/output/{file}")
    
    print("\n🔧 使用方法:")
    print("""
    from src.autoagentsai.slide.pptx2pptx_agent import PPTX2PPTXAgent
    
    # 准备多个数据项
    data_list = [
        {"name": "张三", "score": 95, "photo": "path1.jpg"},
        {"name": "李四", "score": 88, "photo": "path2.jpg"}
    ]
    
    # 使用同一模板页面填入不同内容
    agent = PPTX2PPTXAgent()
    result = agent.fill_multiple_data(
        data_list=data_list,
        template_file_path="template.pptx",
        template_slide_index=0,  # 使用第1页作为模板
        output_file_path="result.pptx"
    )
    """)

if __name__ == "__main__":
    main()
